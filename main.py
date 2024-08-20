import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def fetch_html(url):
    response = requests.get(url)
    response.raise_for_status()
    return response.text

def clean_text(text):
    text = text.replace('_x000D_', '')
    text = text.replace('\r', '')
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def extract_title(html_content):
    soup = BeautifulSoup(html_content, 'lxml')
    title_tag = soup.find('title')
    if title_tag:
        title = title_tag.text.strip()
        title = re.sub(r'[\\/*?:"<>|]', "", title)
        return title
    return "presentation"

def extract_lyrics_and_translation(html_content):
    soup = BeautifulSoup(html_content, 'lxml')

    lyrics = []
    translations = []
    found_lyrics = False
    found_translation = False

    for p in soup.find_all('p'):
        if 'LYRICS' in p.text.upper():
            found_lyrics = True
            continue

        if found_lyrics:
            if 'TRANSLATION' in p.text.upper():
                found_translation = True
                continue

            if found_translation:
                if 'REMARKS' in p.text.upper() or 'CREDITS' in p.text.upper():
                    break
                translations.append(clean_text(p.text.strip()))
            else:
                if 'TRANSLATION' in p.text.upper() or 'REMARKS' in p.text.upper() or 'CREDITS' in p.text.upper():
                    continue
                lyrics.append(clean_text(p.text.strip()))

    if not lyrics:
        raise ValueError("No lyrics found in the document.")

    return lyrics, translations

def split_translations(translations_text):
    # Split translations by numbering pattern (e.g., 1), 2), etc.)
    split_translations = re.split(r'(\d+\)\s)', " ".join(translations_text))
    translations = []
    current_translation = ""
    for segment in split_translations:
        if re.match(r'\d+\)\s', segment):
            if current_translation:
                translations.append(current_translation.strip())
            current_translation = segment
        else:
            current_translation += segment
    if current_translation:
        translations.append(current_translation.strip())
    return translations

def create_ppt(lyrics, translations, output_file, template_file):
    prs = Presentation(template_file)

    verse_pattern = re.compile(r'\(\d+\)')
    verse_indices = [i for i, verse in enumerate(lyrics) if verse_pattern.match(verse)]

    # Adding verses slides
    for i in range(len(verse_indices)):
        start = verse_indices[i]
        end = verse_indices[i + 1] if i + 1 < len(verse_indices) else len(lyrics)
        
        verse = "\n".join(lyrics[start:end]).strip()
        
        # Add the verse slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        text_frame = textbox.text_frame
        text_frame.text = verse

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.color.rgb = RGBColor(0, 0, 0)

        text_frame.word_wrap = True
        textbox.left = Inches(1)
        textbox.top = Inches(1.5)
    
    # Adding all translations slides at the end
    if translations:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
        text_frame = textbox.text_frame

        # Add "Translations" title
        p = text_frame.add_paragraph()
        p.text = "Translations"
        p.font.size = Pt(36)
        #p.alignment = PP_ALIGN.CENTER
        p.font.bold = True

        split_trans = split_translations(translations)
        for translation in split_trans:
            # Add each translation
            p = text_frame.add_paragraph()
            p.text = translation.strip()
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(10)  # Add some space between translations

        text_frame.word_wrap = True
        textbox.left = Inches(1)
        textbox.top = Inches(1.5)

    prs.save(output_file)

def validate_url(url):
    return url.startswith("https://kksongs.org/songs/") and url.endswith(".html")

template_file = "starter.pptx"

while True:
    url = input("Enter the URL of the song: ").strip()
    
    if not validate_url(url):
        print("Invalid URL. Please ensure the URL is from 'https://kksongs.org/songs/' and ends with '.html'.")
        continue
    
    try:
        html_content = fetch_html(url)

        title = extract_title(html_content)

        lyrics, translations = extract_lyrics_and_translation(html_content)

        output_file = f"{title}.pptx"
        create_ppt(lyrics, translations, output_file, template_file)

        print(f"PowerPoint presentation created: {output_file}")
        
        another = input("Do you want to create another presentation? (y/n): ").strip().lower()
        if another != 'y':
            break

    except Exception as e:
        print(f"An error occurred: {e}")
        continue
