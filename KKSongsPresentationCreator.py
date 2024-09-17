import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

class KKSongsPresentationCreator:
    def __init__(self, template_file):
        self.template_file = template_file

    def fetch_html(self, url):
        response = requests.get(url)
        response.raise_for_status()
        return response.text

    @staticmethod
    def clean_text(text):
        text = text.replace('_x000D_', '')
        text = text.replace('', "")
        text = text.replace('\r', '')
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\\x[0-9A-Fa-f]{2}', '', text)
        return text.strip()

    def extract_title(self, html_content):
        soup = BeautifulSoup(html_content, 'lxml')
        title_tag = soup.find('title')
        if title_tag:
            title = title_tag.text.strip()
            title = re.sub(r'[\\/*?:"<>|]', "", title)
            return title
        return "presentation"

    def extract_lyrics_and_translation(self, html_content):
        soup = BeautifulSoup(html_content, 'lxml')
        lyrics = []
        translations = []
        
        found_lyrics = False
        found_translation = False

        for p in soup.find_all('p'):
            # Check for LYRICS section
            if 'LYRICS' in p.text.upper():
                found_lyrics = True
                found_translation = False
                continue

            # Check for TRANSLATION section
            if 'TRANSLATION' in p.text.upper():
                found_translation = True
                found_lyrics = False
                continue

            # Stop at REMARKS or CREDITS section
            if 'REMARKS' in p.text.upper() or 'CREDITS' in p.text.upper():
                break

            # Extract lyrics
            if found_lyrics:
                cleaned_text = self.clean_text(p.text.strip())
                # Add only if it's not empty
                if cleaned_text:
                    lyrics.append(cleaned_text)

            # Extract translations
            elif found_translation:
                cleaned_text = self.clean_text(p.text.strip())
                if cleaned_text:
                    translations.append(cleaned_text)

        # If no lyrics found, raise an error
        if not lyrics:
            raise ValueError("No lyrics found in the document.")

        # Match translations to lyrics by numbering
        numbered_lyrics = self.split_numbered_sections(lyrics)
        numbered_translations = self.split_numbered_sections(translations)

        # Align translations to the corresponding verses
        aligned_lyrics = []
        aligned_translations = []

        for verse_number, verse_lines in numbered_lyrics.items():
            aligned_lyrics.extend(verse_lines)
            # Find the matching translation by number
            if verse_number in numbered_translations:
                aligned_translations.append(" ".join(numbered_translations[verse_number]))
            else:
                aligned_translations.append("")

        return aligned_lyrics, aligned_translations

    def split_numbered_sections(self, text_list):
        """Helper function to split text sections based on numbering (e.g., '1)', '2)')"""
        sections = {}
        current_number = None
        current_lines = []

        for line in text_list:
            # Check if the line starts with a numbering pattern like "1)"
            match = re.match(r'^\(?(\d+)\)?\s*', line)
            if match:
                # If we have a current section, save it
                if current_number is not None:
                    sections[current_number] = current_lines

                # Start a new section
                current_number = match.group(1)
                current_lines = [line]
            else:
                # Add line to the current section
                if current_number is not None:
                    current_lines.append(line)

        # Save the last section
        if current_number is not None:
            sections[current_number] = current_lines

        return sections
    
    def split_numbered_sections(self, text_list):
        """Helper function to split text sections based on numbering (e.g., '1)', '2)')"""
        sections = {}
        current_number = None
        current_lines = []

        for line in text_list:
            # Check if the line starts with a numbering pattern like "1)"
            match = re.match(r'^\(?(\d+)\)?\s*', line)
            if match:
                # If we have a current section, save it
                if current_number is not None:
                    sections[current_number] = current_lines

                # Start a new section
                current_number = match.group(1)
                current_lines = [line]
            else:
                # Add line to the current section
                if current_number is not None:
                    current_lines.append(line)

        # Save the last section
        if current_number is not None:
            sections[current_number] = current_lines

        return sections


    @staticmethod
    def split_translations(translations_text):
        split_translations = re.split(r'(\d+\)\s)', " ".join(translations_text))
        translations = []
        current_translation = ""
        for segment in split_translations:
            if re.match(r'\d+\)\s', segment):
                if current_translation:
                    translations.append(current_translation.strip())
                current_translation = segment
            else:
                current_translation += segment
        if current_translation:
            translations.append(current_translation.strip())
        return translations

    def create_ppt(self, lyrics, translations, output_file):
        prs = Presentation(self.template_file)

        verse_pattern = re.compile(r'\(\d+\)')
        verse_indices = [i for i, verse in enumerate(lyrics) if verse_pattern.match(verse)]

        # Split translations in case they are numbered separately
        split_trans = self.split_translations(translations) if translations else []

        # Adding verses and corresponding translations to slides
        for i in range(len(verse_indices)):
            start = verse_indices[i]
            end = verse_indices[i + 1] if i + 1 < len(verse_indices) else len(lyrics)

            # Join and clean the verse
            verse = "\n".join(lyrics[start:end]).strip()
            translation = split_trans[i] if i < len(split_trans) else ""

            # Split the verse into lines
            verse_lines = verse.split('\n')

            # Get the first line (which might contain the numbering) separately
            first_line = verse_lines[0] if verse_lines[0].startswith('(') else ""
            content_lines = verse_lines[1:] if first_line else verse_lines

            # Determine the number of slides needed
            num_slides = (len(content_lines) + 1) // 2  # Two lines per slide

            for slide_index in range(num_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                verse_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(3))
                text_frame = verse_textbox.text_frame

                # Get the lines for the current slide
                start_line = slide_index * 2
                end_line = min(start_line + 2, len(content_lines))
                slide_lines = content_lines[start_line:end_line]

                # If it's the first slide, include the first line (numbering)
                if slide_index == 0 and first_line:
                    slide_lines.insert(0, first_line)

                verse_text = "\n".join(slide_lines).strip()

                # Adjust font size dynamically based on the line length
                font_size = 32
                max_chars = max(len(line) for line in slide_lines)
                if max_chars > 50:
                    font_size = max(20, int(32 * 50 / max_chars))  # Scale down the font size

                # Add the verse to the slide
                verse_paragraph = text_frame.add_paragraph()
                verse_paragraph.text = verse_text
                verse_paragraph.font.size = Pt(font_size)
                verse_paragraph.alignment = PP_ALIGN.CENTER
                verse_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                text_frame.word_wrap = True

                # Add translation to every sub-slide
                if translation:
                    translation_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(3))
                    trans_text_frame = translation_textbox.text_frame
                    translation_paragraph = trans_text_frame.add_paragraph()
                    translation_paragraph.text = translation.strip()
                    translation_paragraph.font.size = Pt(28)
                    translation_paragraph.alignment = PP_ALIGN.CENTER
                    translation_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    translation_paragraph.space_after = Pt(10)
                    trans_text_frame.word_wrap = True

        prs.save(output_file)

    @staticmethod
    def validate_url(url):
        return url.startswith("https://kksongs.org/songs/") and url.endswith(".html")

